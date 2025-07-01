from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os
import logging
from app.utils import download_image
import tempfile
import requests
from datetime import datetime

# Configuración de logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Paleta de colores corporativa
COLORES = {
    'principal': RGBColor(0x00, 0x3D, 0x7D),    # Azul corporativo oscuro
    'secundario': RGBColor(0x00, 0x84, 0xD1),   # Azul corporativo medio
    'acento': RGBColor(0xFF, 0x8C, 0x00),       # Naranja acento
    'fondo_claro': RGBColor(0xF5, 0xF9, 0xFF),  # Azul muy claro para fondos
    'texto_oscuro': RGBColor(0x2C, 0x3E, 0x50), # Gris azulado para texto
    'blanco': RGBColor(0xFF, 0xFF, 0xFF),       # Blanco puro
    'gris_claro': RGBColor(0xE5, 0xE5, 0xE5)    # Gris claro para bordes
}

# Configuración de fuentes
FUENTES = {
    'titulo': 'Segoe UI',
    'subtitulo': 'Segoe UI Light',
    'cuerpo': 'Segoe UI'
}

def aplicar_estilo_slide(slide):
    """Aplica el estilo base a una diapositiva."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORES['fondo_claro']

def add_footer(slide, text_content):
    """Añade un pie de página mejorado."""
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.3))
    tf = footer.text_frame
    tf.text = text_content
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.name = FUENTES['cuerpo']
    p.font.color.rgb = COLORES['texto_oscuro']
    p.alignment = PP_ALIGN.CENTER

def agregar_logo(slide):
    """Agrega el logo corporativo en la esquina superior izquierda de cada diapositiva."""
    # URL directa al logo de MMI
    logo_url = "https://mmi-e.com/wp-content/uploads/2023/03/LogoMMI_1024x1024.png"
    
    try:
        # Ruta temporal para almacenar el logo descargado
        temp_dir = tempfile.gettempdir()
        logo_path = os.path.join(temp_dir, "logo_mmi.png")
        
        # Intentar descargar el logo si no existe
        if not os.path.exists(logo_path) or os.path.getsize(logo_path) == 0:
            logging.info(f"Descargando logo desde {logo_url}")
            response = requests.get(logo_url, stream=True, timeout=15)
            if response.status_code == 200:
                with open(logo_path, 'wb') as f:
                    for chunk in response.iter_content(chunk_size=8192):
                        f.write(chunk)
                logging.info(f"Logo descargado correctamente a {logo_path}")
            else:
                logging.error(f"Error al descargar logo: {response.status_code}")
                # Intentar con URL alternativa
                alt_logo_url = "https://mmi-e.com/wp-content/uploads/2020/09/logo-mmi.png"
                logging.info(f"Intentando con URL alternativa: {alt_logo_url}")
                response = requests.get(alt_logo_url, stream=True, timeout=15)
                if response.status_code == 200:
                    with open(logo_path, 'wb') as f:
                        for chunk in response.iter_content(chunk_size=8192):
                            f.write(chunk)
                    logging.info(f"Logo alternativo descargado correctamente")
                else:
                    logging.error(f"Error al descargar logo alternativo: {response.status_code}")
                    return
        
        # Verificar que el archivo existe y tiene contenido
        if os.path.exists(logo_path) and os.path.getsize(logo_path) > 0:
            # Posición en esquina superior izquierda
            left = Inches(0.3)  # Cambiado de derecha a izquierda
            top = Inches(0.2)
            width = Inches(1.2)
            
            # Añadir logo con tamaño fijo
            logo = slide.shapes.add_picture(logo_path, left, top, width=width)
            logging.info(f"Logo añadido correctamente a la diapositiva en la esquina izquierda")
        else:
            logging.error(f"El archivo del logo no existe o está vacío: {logo_path}")
            
    except Exception as e:
        logging.error(f"Error al añadir logo: {e}")

def formatear_fecha(fecha_str):
    """
    Convierte el formato de fecha YYYY-MM-DD a DD-MM-YY
    """
    try:
        if not fecha_str or fecha_str == 'N/A':
            return 'N/A'
            
        # Intentar diferentes formatos de fecha posibles
        formatos = ['%Y-%m-%d', '%d-%m-%Y', '%Y/%m/%d', '%d/%m/%Y']
        fecha_obj = None
        
        for formato in formatos:
            try:
                fecha_obj = datetime.strptime(fecha_str, formato)
                break
            except ValueError:
                continue
                
        if not fecha_obj:
            # Si no se pudo parsear, devolver la fecha original
            return fecha_str
            
        # Formatear como DD-MM-YY
        return fecha_obj.strftime('%d-%m-%y')
    except Exception as e:
        logging.error(f"Error al formatear fecha '{fecha_str}': {e}")
        return fecha_str

def formatear_moneda(valor):
    """
    Añade el símbolo de euro (€) a los valores monetarios.
    Preserva los puntos como separadores de miles.
    """
    if not valor or valor == 'N/A':
        return 'N/A'
    
    try:
        # Convertir el valor a string si no lo es
        valor_str = str(valor).strip()
        
        # Si ya tiene el símbolo de euro, lo devolvemos tal cual
        if '€' in valor_str:
            return valor_str
        
        # Añadir el símbolo de euro al final
        return f"{valor_str} €"
    except Exception as e:
        logging.error(f"Error al formatear valor monetario '{valor}': {e}")
        return f"{valor} €"

def crear_portada(pr, datos):
    slide = pr.slides.add_slide(pr.slide_layouts[5])  # Usar layout en blanco
    aplicar_estilo_slide(slide)
    
    # Barra superior decorativa
    top_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(10),
        Inches(1.5)
    )
    top_rect.fill.solid()
    top_rect.fill.fore_color.rgb = COLORES['principal']
    top_rect.line.fill.background()
    
    # Título principal centrado horizontalmente, ajustado para no solaparse con el logo
    title = slide.shapes.add_textbox(
        Inches(1.8),  # Mover hacia la derecha para evitar solapamiento con el logo
        Inches(2.5),
        Inches(6.4),  # Reducir ancho para centrar mejor
        Inches(1)
    )
    tf = title.text_frame
    tf.text = "Informe de Medios"
    p = tf.paragraphs[0]
    p.font.name = FUENTES['titulo']
    p.font.size = Pt(44)
    p.font.color.rgb = COLORES['principal']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo con período
    subtitle = slide.shapes.add_textbox(
        Inches(1.8),  # Ajustar para alinear con el título
        Inches(4),
        Inches(6.4),  # Mismo ancho que el título
        Inches(0.75)
    )
    tf = subtitle.text_frame
    tf.text = f"Período: {datos.get('fechaInicial', 'N/A')} - {datos.get('fechaFinal', 'N/A')}"
    p = tf.paragraphs[0]
    p.font.name = FUENTES['subtitulo']
    p.font.size = Pt(24)
    p.font.color.rgb = COLORES['secundario']
    p.alignment = PP_ALIGN.CENTER
    
    # Añadir logo - asegurar que se añada después de crear todos los elementos
    agregar_logo(slide)
    add_footer(slide, "Informe de Medios")

def crear_metodologia(pr):
    slide = pr.slides.add_slide(pr.slide_layouts[5])  # Usar layout en blanco
    aplicar_estilo_slide(slide)
    
    # Barra de título horizontal
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(10),
        Inches(1)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORES['principal']
    title_box.line.fill.background()
    
    # Título centrado horizontalmente
    title = slide.shapes.add_textbox(
        0,
        Inches(0.2),
        Inches(10),
        Inches(0.6)
    )
    tf = title.text_frame
    tf.text = "Metodología"
    p = tf.paragraphs[0]
    p.font.name = FUENTES['titulo']
    p.font.size = Pt(28)
    p.font.color.rgb = COLORES['blanco']
    p.alignment = PP_ALIGN.CENTER
    
    # Contenido en caja con estilo
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(4)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COLORES['blanco']
    content_box.line.color.rgb = COLORES['gris_claro']
    content_box.shadow.inherit = False
    
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        "• Monitoreo continuo de medios",
        "• Análisis cualitativo y cuantitativo",
        "• Métricas de evaluación:",
        "   - Valor Publicitario Equivalente (VPE)",
        "   - Alcance y audiencia",
        "   - Sentimiento y tono",
        "   - Presencia en diferentes soportes"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = FUENTES['cuerpo']
        p.font.size = Pt(18)
        p.font.color.rgb = COLORES['texto_oscuro']
        if item.startswith("   -"):
            p.level = 1
    
    # Añadir logo
    agregar_logo(slide)
    add_footer(slide, "Metodología - Informe de Medios")

def crear_datos_cobertura(pr, datos, tipo_medio):
    slide = pr.slides.add_slide(pr.slide_layouts[5])
    aplicar_estilo_slide(slide)
    
    medio_data = datos.get(f"{tipo_medio}_raw", {})
    if not medio_data:
        return
    
    # Barra de título horizontal
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(10),
        Inches(1)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORES['principal']
    title_box.line.fill.background()
    
    # Título centrado horizontalmente, ajustado para no solaparse con el logo
    title = slide.shapes.add_textbox(
        Inches(1.8),  # Ajustar para dejar espacio al logo
        Inches(0.2),
        Inches(6.4),  # Reducir ancho para centrar mejor
        Inches(0.6)
    )
    tf = title.text_frame
    tf.text = f"Datos de Cobertura - {tipo_medio}"
    p = tf.paragraphs[0]
    p.font.name = FUENTES['titulo']
    p.font.size = Pt(28)
    p.font.color.rgb = COLORES['blanco']
    p.alignment = PP_ALIGN.CENTER
    
    # Caja de resumen con datos principales
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1.5),
        Inches(8),
        Inches(1.5)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = COLORES['blanco']
    summary_box.line.color.rgb = COLORES['secundario']
    summary_box.shadow.inherit = False
    
    tf = summary_box.text_frame
    tf.word_wrap = True
    
    # Formatear valores monetarios con símbolo de euro
    vpe_formateado = formatear_moneda(medio_data.get('total_vpe', 'N/A'))
    vc_formateado = formatear_moneda(medio_data.get('total_vc', 'N/A'))
    
    # Datos de resumen con iconos
    p = tf.add_paragraph()
    p.text = f"📊 Total de Noticias: {medio_data.get('cantidad_noticias', 'N/A')}"
    p.font.name = FUENTES['cuerpo']
    p.font.size = Pt(14)
    p.font.color.rgb = COLORES['texto_oscuro']
    
    p = tf.add_paragraph()
    p.text = f"👥 Audiencia Total: {medio_data.get('total_audiencia', 'N/A')}"
    p.font.name = FUENTES['cuerpo']
    p.font.size = Pt(14)
    p.font.color.rgb = COLORES['texto_oscuro']
    
    p = tf.add_paragraph()
    p.text = f"💰 VPE: {vpe_formateado} | VC: {vc_formateado}"
    p.font.name = FUENTES['cuerpo']
    p.font.size = Pt(14)
    p.font.color.rgb = COLORES['texto_oscuro']
    
    # Lista de noticias
    noticias_list = medio_data.get("noticias", [])
    if noticias_list:
        # Calcular cuántas noticias podemos mostrar por diapositiva (máx 4)
        max_noticias_por_slide = min(4, len(noticias_list))
        
        news_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(3.2),
            Inches(8),
            Inches(3)
        )
        news_box.fill.solid()
        news_box.fill.fore_color.rgb = COLORES['blanco']
        news_box.line.color.rgb = COLORES['gris_claro']
        news_box.shadow.inherit = False
        
        tf = news_box.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = "📰 Noticias Destacadas"
        p.font.name = FUENTES['subtitulo']
        p.font.size = Pt(16)
        p.font.color.rgb = COLORES['secundario']
        p.space_after = Pt(10)
        
        # Añadir las primeras noticias
        for i in range(max_noticias_por_slide):
            if i >= len(noticias_list):
                break
                
            noticia = noticias_list[i]
            
            # Formatear la fecha a DD-MM-AA
            fecha_formateada = formatear_fecha(noticia.get('fecha', 'N/A'))
            
            p = tf.add_paragraph()
            p.text = f"📅 {fecha_formateada} - {noticia.get('titulo', 'N/A')}"
            p.font.name = FUENTES['cuerpo']
            p.font.size = Pt(12)
            p.font.color.rgb = COLORES['texto_oscuro']
            p.font.bold = True
            p.space_before = Pt(5)
            p.space_after = Pt(2)
            
            # Párrafo con hipervínculo
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"     {noticia.get('titular', 'N/A')}"
            run.font.name = FUENTES['cuerpo']
            run.font.size = Pt(11)
            run.font.color.rgb = COLORES['secundario']
            run.font.italic = True
            
            # Añadir hipervínculo al titular
            if 'url' in noticia and noticia['url']:
                hlink = run.hyperlink
                hlink.address = noticia['url']
            elif 'link' in noticia and noticia['link']:
                hlink = run.hyperlink
                hlink.address = noticia['link']
            
            p.space_after = Pt(8)
        
        # Si hay más noticias, crear una nueva diapositiva
        if len(noticias_list) > max_noticias_por_slide:
            slide_continuacion = pr.slides.add_slide(pr.slide_layouts[5])
            aplicar_estilo_slide(slide_continuacion)
            
            # Barra de título
            title_box = slide_continuacion.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                Inches(10),
                Inches(1)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = COLORES['principal']
            title_box.line.fill.background()
            
            # Título principal
            title = slide_continuacion.shapes.add_textbox(
                Inches(1.8),
                Inches(0.2),
                Inches(6.4),
                Inches(0.6)
            )
            tf = title.text_frame
            tf.text = f"Datos de Cobertura - {tipo_medio}"
            p = tf.paragraphs[0]
            p.font.name = FUENTES['titulo']
            p.font.size = Pt(28)
            p.font.color.rgb = COLORES['blanco']
            p.alignment = PP_ALIGN.CENTER
            
            # Subtítulo (Continuación) debajo del título principal
            subtitle = slide_continuacion.shapes.add_textbox(
                Inches(1.8),
                Inches(0.6),  # Justo debajo del título
                Inches(6.4),
                Inches(0.4)
            )
            tf = subtitle.text_frame
            tf.text = "(Continuación)"
            p = tf.paragraphs[0]
            p.font.name = FUENTES['subtitulo']
            p.font.size = Pt(16)  # 40% más pequeño que el título principal
            p.font.color.rgb = COLORES['blanco']
            p.alignment = PP_ALIGN.CENTER
            
            # Caja de noticias
            news_box = slide_continuacion.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1),
                Inches(1.5),
                Inches(8),
                Inches(5)
            )
            news_box.fill.solid()
            news_box.fill.fore_color.rgb = COLORES['blanco']
            news_box.line.color.rgb = COLORES['gris_claro']
            news_box.shadow.inherit = False
            
            tf = news_box.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = "📰 Noticias Destacadas"
            p.font.name = FUENTES['subtitulo']
            p.font.size = Pt(16)
            p.font.color.rgb = COLORES['secundario']
            p.space_after = Pt(10)
            
            # Añadir las noticias restantes
            for i in range(max_noticias_por_slide, len(noticias_list)):
                noticia = noticias_list[i]
                
                # Formatear la fecha a DD-MM-AA
                fecha_formateada = formatear_fecha(noticia.get('fecha', 'N/A'))
                
                p = tf.add_paragraph()
                p.text = f"📅 {fecha_formateada} - {noticia.get('titulo', 'N/A')}"
                p.font.name = FUENTES['cuerpo']
                p.font.size = Pt(12)
                p.font.color.rgb = COLORES['texto_oscuro']
                p.font.bold = True
                p.space_before = Pt(5)
                p.space_after = Pt(2)
                
                # Párrafo con hipervínculo
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"     {noticia.get('titular', 'N/A')}"
                run.font.name = FUENTES['cuerpo']
                run.font.size = Pt(11)
                run.font.color.rgb = COLORES['secundario']
                run.font.italic = True
                
                # Añadir hipervínculo al titular
                if 'url' in noticia and noticia['url']:
                    hlink = run.hyperlink
                    hlink.address = noticia['url']
                elif 'link' in noticia and noticia['link']:
                    hlink = run.hyperlink
                    hlink.address = noticia['link']
                
                p.space_after = Pt(8)
            
            # Añadir logo a la diapositiva de continuación
            agregar_logo(slide_continuacion)
            add_footer(slide_continuacion, f"Cobertura {tipo_medio} - Informe de Medios")
    
    # Añadir logo
    agregar_logo(slide)
    add_footer(slide, f"Cobertura {tipo_medio} - Informe de Medios")

def crear_graficos(pr, datos):
    urls = datos.get("urls", [])
    if not urls:
        return
    
    # Ordenar los gráficos según el tipo y la secuencia deseada
    graficos_ordenados = {
        'general': [],     # Gráficos generales (barras, tortas)
        'prensa': [],      # Gráficos de prensa
        'radio': [],       # Gráficos de radio
        'tv': [],          # Gráficos de TV
        'digitales': []    # Gráficos de medios digitales
    }
    
    # Clasificar los gráficos según su tipo
    for url in urls:
        if "top10_vpe_prensa" in url:
            graficos_ordenados['prensa'].append(url)
        elif "top10_vpe_radio" in url:
            graficos_ordenados['radio'].append(url)
        elif "top10_vpe_tv" in url:
            graficos_ordenados['tv'].append(url)
        elif "top10_vpe_medios_digitales" in url or "top10_vpe_digitales" in url:
            graficos_ordenados['digitales'].append(url)
        else:
            # Gráficos generales (barras, tortas)
            graficos_ordenados['general'].append(url)
    
    # Procesar primero los gráficos generales
    for url in graficos_ordenados['general']:
        crear_diapositiva_grafico(pr, url)
    
    # Procesar los gráficos en el orden específico: prensa, radio, TV, digitales
    for tipo in ['prensa', 'radio', 'tv', 'digitales']:
        for url in graficos_ordenados[tipo]:
            crear_diapositiva_grafico(pr, url)

def crear_diapositiva_grafico(pr, url):
    """Crea una diapositiva para un gráfico específico."""
    slide = pr.slides.add_slide(pr.slide_layouts[5])
    aplicar_estilo_slide(slide)
    
    # Determinar tipo de gráfico y título
    tipo_grafico = ""
    subtitulo = ""
    
    if "vpe_barra" in url:
        tipo_grafico = "VPE por Medio"
    elif "vpe_torta" in url:
        tipo_grafico = "Distribución de VPE"
    elif "impactos_barra" in url:
        tipo_grafico = "Impactos por Medio"
    elif "impactos_torta" in url:
        tipo_grafico = "Distribución de Impactos"
    elif "top10_vpe_prensa" in url:
        tipo_grafico = "Top VPE - Prensa"
        subtitulo = "Detalles"
    elif "top10_vpe_radio" in url:
        tipo_grafico = "Top VPE - Radio"
        subtitulo = "Detalles"
    elif "top10_vpe_tv" in url:
        tipo_grafico = "Top VPE - TV"
        subtitulo = "Detalles"
    elif "top10_vpe_medios_digitales" in url or "top10_vpe_digitales" in url:
        tipo_grafico = "Top VPE - Medios Digitales"
        subtitulo = "Detalles"
    else:
        # Extraer nombre del medio si es un top10 específico
        if "top10_vpe_" in url:
            medio = url.split("top10_vpe_")[1].split(".")[0].replace("_", " ").title()
            tipo_grafico = f"Top VPE - {medio}"
            subtitulo = "Detalles"
    
    # Barra de título
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(10),
        Inches(1)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORES['principal']
    title_box.line.fill.background()
    
    # Título centrado horizontalmente, ajustado para no solaparse con el logo
    title = slide.shapes.add_textbox(
        Inches(1.8),  # Ajustar para dejar espacio al logo
        Inches(0.2),
        Inches(6.4),  # Reducir ancho para centrar mejor
        Inches(0.6)
    )
    tf = title.text_frame
    tf.text = tipo_grafico
    p = tf.paragraphs[0]
    p.font.name = FUENTES['titulo']
    p.font.size = Pt(28)
    p.font.color.rgb = COLORES['blanco']
    p.alignment = PP_ALIGN.CENTER
    
    # Añadir subtítulo si es necesario
    if subtitulo:
        subtitle = slide.shapes.add_textbox(
            Inches(0),
            Inches(1.2),
            Inches(10),
            Inches(0.4)
        )
        tf = subtitle.text_frame
        tf.text = subtitulo
        p = tf.paragraphs[0]
        p.font.name = FUENTES['subtitulo']
        p.font.size = Pt(20)
        p.font.color.rgb = COLORES['secundario']
        p.alignment = PP_ALIGN.CENTER
    
    # Primero descargamos la imagen del gráfico
    img_path = None
    try:
        # Usar requests para descargar la imagen
        temp_dir = tempfile.gettempdir()
        img_path = os.path.join(temp_dir, f"grafico_{hash(url)}.png")
        
        response = requests.get(url, stream=True, timeout=15)
        if response.status_code == 200:
            with open(img_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            logging.info(f"Gráfico descargado correctamente: {url}")
        else:
            logging.error(f"Error al descargar gráfico: {response.status_code}")
            img_path = None
    except Exception as e:
        logging.error(f"Error descargando gráfico {url}: {e}")
        img_path = None
    
    # Área de contenido principal (centrado en la diapositiva)
    content_area_top = Inches(1.7) if subtitulo else Inches(1.5)
    content_area_height = Inches(4.8)
    content_area_left = Inches(0.5)
    content_area_width = Inches(9)
    
    # Añadir marco decorativo para el gráfico
    chart_frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        content_area_left,
        content_area_top,
        content_area_width,
        content_area_height
    )
    chart_frame.fill.solid()
    chart_frame.fill.fore_color.rgb = COLORES['blanco']
    chart_frame.line.color.rgb = COLORES['secundario']
    chart_frame.line.width = Pt(2)
    chart_frame.shadow.inherit = False
    
    # Intentar insertar el gráfico si se descargó correctamente
    if img_path and os.path.exists(img_path) and os.path.getsize(img_path) > 0:
        try:
            # Obtener dimensiones de la imagen
            img = Image.open(img_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height
            
            # Calcular tamaño manteniendo proporción
            # Dejamos margen de 0.5 pulgadas a cada lado
            max_width = content_area_width - Inches(1)
            max_height = content_area_height - Inches(0.5)
            
            # Calcular dimensiones para mantener proporción y centrar
            if aspect_ratio > 1:  # Imagen más ancha que alta
                target_width = max_width
                target_height = target_width / aspect_ratio
                if target_height > max_height:
                    target_height = max_height
                    target_width = target_height * aspect_ratio
            else:  # Imagen más alta que ancha
                target_height = max_height
                target_width = target_height * aspect_ratio
                if target_width > max_width:
                    target_width = max_width
                    target_height = target_width / aspect_ratio
            
            # Calcular posición para centrar perfectamente
            left = content_area_left + (content_area_width - target_width) / 2
            top = content_area_top + (content_area_height - target_height) / 2
            
            # Insertar imagen del gráfico
            pic = slide.shapes.add_picture(
                img_path,
                left,
                top,
                width=target_width,
                height=target_height
            )
            
            # Asegurarse de que el gráfico esté en primer plano
            pic.z_order = -1  # Poner en primer plano
            
        except Exception as e:
            logging.error(f"Error al insertar gráfico {url}: {e}")
            
            # Mostrar mensaje de error
            error_box = slide.shapes.add_textbox(
                content_area_left + Inches(1.5),
                content_area_top + Inches(2),
                Inches(6),
                Inches(1)
            )
            tf = error_box.text_frame
            tf.text = "Error al cargar el gráfico"
            p = tf.paragraphs[0]
            p.font.color.rgb = COLORES['acento']
            p.font.size = Pt(16)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        finally:
            # Limpiar archivo temporal
            try:
                if os.path.exists(img_path):
                    os.remove(img_path)
            except:
                pass
    else:
        # Mostrar mensaje de error si no se pudo descargar
        error_box = slide.shapes.add_textbox(
            content_area_left + Inches(1.5),
            content_area_top + Inches(2),
            Inches(6),
            Inches(1)
        )
        tf = error_box.text_frame
        tf.text = "Error al descargar el gráfico"
        p = tf.paragraphs[0]
        p.font.color.rgb = COLORES['acento']
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    # Añadir logo
    agregar_logo(slide)
    add_footer(slide, f"{tipo_grafico} - Informe de Medios")

def crear_vpe_totales(pr, datos):
    slide = pr.slides.add_slide(pr.slide_layouts[5])
    aplicar_estilo_slide(slide)
    
    # Barra de título horizontal
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(10),
        Inches(1)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = COLORES['principal']
    title_box.line.fill.background()
    
    # Título centrado horizontalmente, ajustado para no solaparse con el logo
    title = slide.shapes.add_textbox(
        Inches(1.8),  # Ajustar para dejar espacio al logo
        Inches(0.2),
        Inches(6.4),  # Reducir ancho para centrar mejor
        Inches(0.6)
    )
    tf = title.text_frame
    tf.text = "Valor Publicitario Equivalente (VPE) Total"
    p = tf.paragraphs[0]
    p.font.name = FUENTES['titulo']
    p.font.size = Pt(28)
    p.font.color.rgb = COLORES['blanco']
    p.alignment = PP_ALIGN.CENTER
    
    # Caja de datos VPE
    vpe_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(3.5)
    )
    vpe_box.fill.solid()
    vpe_box.fill.fore_color.rgb = COLORES['blanco']
    vpe_box.line.color.rgb = COLORES['secundario']
    vpe_box.shadow.inherit = False
    
    tf = vpe_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Resumen de Valor Publicitario"
    p.font.name = FUENTES['subtitulo']
    p.font.size = Pt(20)
    p.font.color.rgb = COLORES['secundario']
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)
    
    # Formatear el valor VPE con símbolo de euro
    vpe_formateado = formatear_moneda(datos.get('totalGlobalVPE', 'N/A'))
    
    p = tf.add_paragraph()
    p.text = f"VPE Total: {vpe_formateado}"
    p.font.name = FUENTES['cuerpo']
    p.font.size = Pt(28)
    p.font.color.rgb = COLORES['principal']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Añadir logo
    agregar_logo(slide)
    add_footer(slide, "VPE Total - Informe de Medios")

def generar_pptx(data, filename):
    pr = Presentation()
    
    # Validación de datos de entrada
    if not isinstance(data, (list, dict)):
        logging.error("Input data must be a list or dict.")
        raise ValueError("Input data must be a list or dict.")
    
    datos = data[0] if isinstance(data, list) and data else data
    if not isinstance(datos, dict):
        logging.error("No se pudo extraer el objeto de datos principal.")
        raise ValueError("No se pudo extraer el objeto de datos principal.")
    
    # Generar estructura de presentación
    crear_portada(pr, datos)
    
    # Datos de cobertura por tipo de medio
    for medio in ["TV", "Radio", "Prensa", "Medios Digitales"]:
        crear_datos_cobertura(pr, datos, medio)
    
    crear_vpe_totales(pr, datos)
    crear_graficos(pr, datos)
    
    # Guardar presentación
    output_path = f"/tmp/{filename}"
    try:
        pr.save(output_path)
        logging.info(f"Presentación PPTX guardada en: {output_path}")
    except Exception as e:
        logging.error(f"Error al guardar la presentación PPTX: {e}")
        raise
    
    return output_path
